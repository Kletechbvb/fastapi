# from fastapi import APIRouter, UploadFile, Form
# from pymongo import MongoClient
# from bson import ObjectId
# from datetime import datetime
# import docx, io, asyncio, httpx
# from pdf2image import convert_from_bytes
# from pptx import Presentation        # NEW for pptx
# import openpyxl                      # NEW for xlsx
# import xlrd                          # NEW for xls

# # MongoDB setup (hardcoded)
# MONGO_URI = "mongodb+srv://chatpdfxai_db_user:esfmQRoJQZpJ7if3@cluster0.xzatb0d.mongodb.net/chatpdf?retryWrites=true&w=majority&appName=Cluster0"
# client = MongoClient(MONGO_URI)
# db = client["chatpdf"]
# chats_collection = db["chats"]

# # OCR.space config (hardcoded)
# OCR_API_KEY = "K85634264488957"
# OCR_URL = "https://api.ocr.space/parse/image"

# router = APIRouter(prefix="/chat", tags=["chat"])


# async def ocr_image_page(session, buf: io.BytesIO, page_num: int) -> str:
#     """Send a single page image to OCR.space"""
#     files = {"file": (f"page_{page_num}.png", buf, "image/png")}
#     data = {"apikey": OCR_API_KEY, "language": "eng"}

#     resp = await session.post(OCR_URL, files=files, data=data)
#     result = resp.json()

#     if result.get("IsErroredOnProcessing"):
#         return f"\n❌ OCR failed on page {page_num}: {result.get('ErrorMessage', 'Unknown error')}\n"

#     text = ""
#     for parsed_result in result.get("ParsedResults", []):
#         text += parsed_result.get("ParsedText", "") + "\n"
#     return text


# async def extract_text_from_file(file: UploadFile) -> tuple[str, int]:
#     """Extract text from uploaded file, return (text, page_count)."""
#     filename = file.filename.lower()

#     if filename.endswith(".docx"):
#         d = docx.Document(file.file)
#         text = "\n".join([p.text for p in d.paragraphs])
#         return text, 1

#     elif filename.endswith(".txt"):
#         try:
#             text = file.file.read().decode("utf-8")
#         except UnicodeDecodeError:
#             file.file.seek(0)
#             text = file.file.read().decode("latin-1")
#         return text, 1

#     elif filename.endswith(".pdf"):
#         file.file.seek(0)
#         pdf_bytes = file.file.read()

#         # Convert each PDF page into an image
#         pages = convert_from_bytes(pdf_bytes, dpi=200)
#         tasks = []

#         async with httpx.AsyncClient(timeout=None) as session:
#             for idx, page in enumerate(pages, start=1):
#                 buf = io.BytesIO()
#                 page.save(buf, format="PNG")
#                 buf.seek(0)
#                 tasks.append(ocr_image_page(session, buf, idx))

#             results = await asyncio.gather(*tasks)

#         text = "\n".join(results).strip()
#         return text, len(pages)

#     elif filename.endswith((".png", ".jpg", ".jpeg")):
#         file.file.seek(0)
#         async with httpx.AsyncClient(timeout=None) as session:
#             files = {"file": (file.filename, file.file, file.content_type)}
#             data = {"apikey": OCR_API_KEY, "language": "eng"}
#             resp = await session.post(OCR_URL, files=files, data=data)
#             result = resp.json()

#             if result.get("IsErroredOnProcessing"):
#                 return f"❌ OCR failed: {result.get('ErrorMessage', 'Unknown error')}", 1

#             text = ""
#             for parsed_result in result.get("ParsedResults", []):
#                 text += parsed_result.get("ParsedText", "") + "\n"
#             return text.strip(), 1

#     # 🔹 PowerPoint support (slides, tables, notes)
#     elif filename.endswith(".pptx"):
#         file.file.seek(0)
#         prs = Presentation(file.file)
#         text = ""

#         for idx, slide in enumerate(prs.slides, start=1):
#             text += f"\n--- Slide {idx} ---\n"
#             for shape in slide.shapes:
#                 if hasattr(shape, "text") and shape.text.strip():
#                     text += shape.text + "\n"

#                 # Handle tables
#                 if shape.has_table:
#                     for row in shape.table.rows:
#                         row_text = " | ".join([cell.text.strip() for cell in row.cells])
#                         text += row_text + "\n"

#             # Handle speaker notes
#             if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
#                 notes = slide.notes_slide.notes_text_frame.text.strip()
#                 if notes:
#                     text += f"Notes: {notes}\n"

#         return text.strip(), len(prs.slides)

#     # 🔹 Excel support (.xlsx)
#     elif filename.endswith(".xlsx"):
#         file.file.seek(0)
#         wb = openpyxl.load_workbook(file.file, read_only=True)
#         text = ""
#         for sheet in wb.sheetnames:
#             text += f"\n--- Sheet: {sheet} ---\n"
#             ws = wb[sheet]
#             for row in ws.iter_rows(values_only=True):
#                 text += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
#         return text.strip(), len(wb.sheetnames)

#     # 🔹 Excel support (.xls)
#     elif filename.endswith(".xls"):
#         file.file.seek(0)
#         wb = xlrd.open_workbook(file_contents=file.file.read())
#         text = ""
#         for sheet in wb.sheets():
#             text += f"\n--- Sheet: {sheet.name} ---\n"
#             for row_idx in range(sheet.nrows):
#                 row = sheet.row(row_idx)
#                 text += " | ".join([str(cell.value) for cell in row]) + "\n"
#         return text.strip(), wb.nsheets

#     else:
#         return "Unsupported file format.", 0


# @router.post("/create")
# async def create_chat(user_email: str = Form(...), file: UploadFile = Form(...)):
#     extracted_text, page_count = await extract_text_from_file(file)

#     chat_doc = {
#         "user_email": user_email,
#         "file_name": file.filename,
#         "context": extracted_text,
#         "messages": [],
#         "created_at": datetime.utcnow().isoformat()
#     }
#     result = chats_collection.insert_one(chat_doc)

#     return {
#         "status": "success",
#         "chat_id": str(result.inserted_id),
#         "page_count": page_count,
#         "text_length": len(extracted_text),
#         "pages_snippet": extracted_text[:1000]  # bigger preview for debugging
#     }


# @router.get("/history/{chat_id}")
# def get_chat_history(chat_id: str):
#     chat = chats_collection.find_one({"_id": ObjectId(chat_id)}, {"messages": 1, "file_name": 1})
#     if not chat:
#         return {"status": "error", "message": "Chat not found ❌"}
#     return {
#         "status": "success",
#         "file_name": chat["file_name"],
#         "messages": chat.get("messages", [])
#     }


# @router.get("/getall")
# def get_all_chats(user_email: str):
#     chats = chats_collection.find({"user_email": user_email}, {"file_name": 1, "created_at": 1})
#     result = [
#         {"chat_id": str(c["_id"]), "file_name": c["file_name"], "created_at": c.get("created_at")}
#         for c in chats
#     ]
#     return {"status": "success", "chats": result}


# @router.delete("/delete/{chat_id}")
# def delete_chat(chat_id: str, user_email: str):
#     chat = chats_collection.find_one({"_id": ObjectId(chat_id)})
#     if not chat:
#         return {"status": "error", "message": "Chat not found ❌"}

#     if chat["user_email"] != user_email:
#         return {"status": "error", "message": "Not authorized to delete this chat ❌"}

#     chats_collection.delete_one({"_id": ObjectId(chat_id)})
#     return {"status": "success", "message": f"Chat {chat_id} deleted ✅"}



from fastapi import FastAPI, APIRouter, UploadFile, Form
from pymongo import MongoClient
from bson import ObjectId
from datetime import datetime
import docx, io, asyncio
from pdf2image import convert_from_bytes
from pptx import Presentation
import openpyxl
import xlrd
import pytesseract
from PIL import Image

# ==========================================================
# 🧩 Setup
# ==========================================================

app = FastAPI(title="ChatPDFXAI", version="2.0")

router = APIRouter(prefix="/chat", tags=["chat"])

# MongoDB setup (hardcoded for simplicity — replace with env var in production)
MONGO_URI = "mongodb+srv://chatpdfxai_db_user:esfmQRoJQZpJ7if3@cluster0.xzatb0d.mongodb.net/chatpdf?retryWrites=true&w=majority&appName=Cluster0"
client = MongoClient(MONGO_URI)
db = client["chatpdf"]
chats_collection = db["chats"]


# ==========================================================
# 🔹 Helper: Local OCR using Tesseract
# ==========================================================

def ocr_image(image: Image.Image) -> str:
    """Perform OCR using Tesseract on a PIL image."""
    text = pytesseract.image_to_string(image, lang="eng")
    return text.strip()


# ==========================================================
# 🧠 Text Extraction from Uploaded Files
# ==========================================================

async def extract_text_from_file(file: UploadFile) -> tuple[str, int]:
    """Extract text from uploaded file using local Tesseract OCR (no API limits)."""
    filename = file.filename.lower()

    if filename.endswith(".docx"):
        d = docx.Document(file.file)
        text = "\n".join([p.text for p in d.paragraphs])
        return text, 1

    elif filename.endswith(".txt"):
        try:
            text = file.file.read().decode("utf-8")
        except UnicodeDecodeError:
            file.file.seek(0)
            text = file.file.read().decode("latin-1")
        return text, 1

    elif filename.endswith(".pdf"):
        file.file.seek(0)
        pdf_bytes = file.file.read()

        pages = convert_from_bytes(pdf_bytes, dpi=150)
        text = ""

        for idx, page in enumerate(pages, start=1):
            page_text = ocr_image(page)
            text += f"\n--- Page {idx} ---\n{page_text}\n"

        return text.strip(), len(pages)

    elif filename.endswith((".png", ".jpg", ".jpeg")):
        file.file.seek(0)
        image = Image.open(file.file)
        text = ocr_image(image)
        return text, 1

    elif filename.endswith(".pptx"):
        file.file.seek(0)
        prs = Presentation(file.file)
        text = ""
        for idx, slide in enumerate(prs.slides, start=1):
            text += f"\n--- Slide {idx} ---\n"
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    text += shape.text + "\n"
                if shape.has_table:
                    for row in shape.table.rows:
                        text += " | ".join([cell.text.strip() for cell in row.cells]) + "\n"
            if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
                notes = slide.notes_slide.notes_text_frame.text.strip()
                if notes:
                    text += f"Notes: {notes}\n"
        return text.strip(), len(prs.slides)

    elif filename.endswith(".xlsx"):
        file.file.seek(0)
        wb = openpyxl.load_workbook(file.file, read_only=True)
        text = ""
        for sheet in wb.sheetnames:
            text += f"\n--- Sheet: {sheet} ---\n"
            ws = wb[sheet]
            for row in ws.iter_rows(values_only=True):
                text += " | ".join([str(cell) if cell else "" for cell in row]) + "\n"
        return text.strip(), len(wb.sheetnames)

    elif filename.endswith(".xls"):
        file.file.seek(0)
        wb = xlrd.open_workbook(file_contents=file.file.read())
        text = ""
        for sheet in wb.sheets():
            text += f"\n--- Sheet: {sheet.name} ---\n"
            for row_idx in range(sheet.nrows):
                row = sheet.row(row_idx)
                text += " | ".join([str(cell.value) for cell in row]) + "\n"
        return text.strip(), wb.nsheets

    else:
        return "Unsupported file format.", 0


# ==========================================================
# 🧩 Routes
# ==========================================================

@router.post("/create")
async def create_chat(user_email: str = Form(...), file: UploadFile = Form(...)):
    extracted_text, page_count = await extract_text_from_file(file)

    chat_doc = {
        "user_email": user_email,
        "file_name": file.filename,
        "context": extracted_text,
        "messages": [],
        "created_at": datetime.utcnow().isoformat()
    }

    result = chats_collection.insert_one(chat_doc)

    return {
        "status": "success",
        "chat_id": str(result.inserted_id),
        "page_count": page_count,
        "text_length": len(extracted_text),
        "pages_snippet": extracted_text[:1000],
    }


@router.get("/history/{chat_id}")
def get_chat_history(chat_id: str):
    chat = chats_collection.find_one({"_id": ObjectId(chat_id)}, {"messages": 1, "file_name": 1})
    if not chat:
        return {"status": "error", "message": "Chat not found ❌"}
    return {
        "status": "success",
        "file_name": chat["file_name"],
        "messages": chat.get("messages", []),
    }


@router.get("/getall")
def get_all_chats(user_email: str):
    chats = chats_collection.find({"user_email": user_email}, {"file_name": 1, "created_at": 1})
    result = [
        {"chat_id": str(c["_id"]), "file_name": c["file_name"], "created_at": c.get("created_at")}
        for c in chats
    ]
    return {"status": "success", "chats": result}


@router.delete("/delete/{chat_id}")
def delete_chat(chat_id: str, user_email: str):
    chat = chats_collection.find_one({"_id": ObjectId(chat_id)})
    if not chat:
        return {"status": "error", "message": "Chat not found ❌"}
    if chat["user_email"] != user_email:
        return {"status": "error", "message": "Not authorized ❌"}
    chats_collection.delete_one({"_id": ObjectId(chat_id)})
    return {"status": "success", "message": f"Chat {chat_id} deleted ✅"}


# ==========================================================
# 🔹 Register Router
# ==========================================================

app.include_router(router)


@app.get("/")
def root():
    return {"status": "running ✅", "message": "ChatPDFXAI OCR (local Tesseract version)"}
